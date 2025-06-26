from typing import List
import os
from pptx import Presentation
from PIL import Image
import pytesseract
import io
from langchain_core.documents import Document
import pandas as pd

class PowerPointLoader:
    def __init__(self, file_path: str):
        """Initialise le loader avec le chemin du fichier PowerPoint."""
        self.file_path = file_path
        self.presentation = Presentation(file_path)

    def _extract_table_content(self, shape) -> str:
        """Extrait le contenu d'un tableau."""
        try:
            if not hasattr(shape, "table"):
                return ""
                
            table_data = []
            for row in shape.table.rows:
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text.strip())
                table_data.append(row_data)
            
            if not table_data:
                return ""
            
            # Convertir en DataFrame pour une meilleure lisibilité
            df = pd.DataFrame(table_data[1:], columns=table_data[0] if table_data else None)
            return f"Tableau:\n{df.to_string()}"
        except Exception as e:
            print(f"Note: Impossible d'extraire le tableau: {str(e)}")
            return ""

    def _extract_image_text(self, shape) -> str:
        """Extrait le texte d'une image en utilisant OCR."""
        try:
            if not hasattr(shape, "image"):
                return ""
                
            image_stream = io.BytesIO()
            image = shape.image
            image.save(image_stream)
            image_stream.seek(0)
            
            # Ouvrir l'image avec Pillow
            img = Image.open(image_stream)
            
            # Extraire le texte avec Tesseract
            text = pytesseract.image_to_string(img, lang='fra')
            return f"Texte extrait de l'image:\n{text.strip()}" if text.strip() else ""
        except Exception as e:
            print(f"Note: Impossible d'extraire le texte de l'image: {str(e)}")
            return ""

    def _process_shape(self, shape) -> str:
        """Traite un shape et retourne son contenu."""
        content = []
        
        # Extraire le texte si disponible
        if hasattr(shape, "text") and shape.text.strip():
            content.append(shape.text.strip())
            
        # Extraire le contenu du tableau si disponible
        table_content = self._extract_table_content(shape)
        if table_content:
            content.append(table_content)
            
        # Extraire le texte de l'image si disponible
        image_content = self._extract_image_text(shape)
        if image_content:
            content.append(image_content)
            
        return "\n".join(content)

    def load(self) -> List[Document]:
        """Charge la présentation et retourne une liste de Documents."""
        documents = []
        
        for slide_number, slide in enumerate(self.presentation.slides, 1):
            # Extraire le titre de la slide
            title = ""
            if slide.shapes.title:
                title = slide.shapes.title.text.strip()

            # Initialiser les conteneurs pour chaque type de contenu
            contents = []

            # Traiter chaque shape dans la slide
            for shape in slide.shapes:
                if shape == slide.shapes.title:  # Éviter de dupliquer le titre
                    continue
                    
                content = self._process_shape(shape)
                if content:
                    contents.append(content)

            # Construire le contenu structuré
            content_parts = [
                f"Numéro de slide: {slide_number}",
                f"Titre: {title}" if title else "Titre: Non défini",
                *contents
            ]

            # Créer le document avec les métadonnées
            doc = Document(
                page_content="\n".join(part for part in content_parts if part),
                metadata={
                    "source": os.path.basename(self.file_path),
                    "slide_number": slide_number,
                    "slide_title": title,
                    "type": "powerpoint"
                }
            )
            documents.append(doc)

        return documents 